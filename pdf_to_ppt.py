"""
PDF to PowerPoint Converter
Converts each page of a PDF to an image and adds it to a PowerPoint presentation.
"""

import fitz  # PyMuPDF
from pptx import Presentation
from pptx.util import Inches
import os
import io

def pdf_to_ppt(pdf_path, output_path):
    """Convert PDF to PowerPoint with each page as a slide."""
    
    # Open the PDF
    pdf_document = fitz.open(pdf_path)
    
    # Create a PowerPoint presentation (16:9 widescreen)
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    
    # Blank layout
    blank_layout = prs.slide_layouts[6]
    
    print(f"Converting {pdf_document.page_count} pages...")
    
    for page_num in range(pdf_document.page_count):
        # Get the page
        page = pdf_document[page_num]
        
        # Render page to image (high resolution)
        mat = fitz.Matrix(2.5, 2.5)  # 2.5x zoom for better quality
        pix = page.get_pixmap(matrix=mat)
        
        # Convert to PNG bytes
        img_bytes = pix.tobytes("png")
        
        # Create a slide
        slide = prs.slides.add_slide(blank_layout)
        
        # Save image temporarily
        temp_img_path = f"temp_slide_{page_num}.png"
        with open(temp_img_path, "wb") as f:
            f.write(img_bytes)
        
        # Add image to slide (full slide)
        slide.shapes.add_picture(
            temp_img_path,
            Inches(0),
            Inches(0),
            width=prs.slide_width,
            height=prs.slide_height
        )
        
        # Remove temp file
        os.remove(temp_img_path)
        
        print(f"  Page {page_num + 1} done")
    
    # Save the PowerPoint
    prs.save(output_path)
    pdf_document.close()
    
    print(f"\n✅ PowerPoint saved to: {output_path}")

if __name__ == "__main__":
    pdf_path = r"C:\Users\gkaru\Downloads\zap.pdf"
    output_path = r"C:\Users\gkaru\Downloads\zap.pptx"
    
    if os.path.exists(pdf_path):
        pdf_to_ppt(pdf_path, output_path)
    else:
        print(f"❌ PDF not found: {pdf_path}")
